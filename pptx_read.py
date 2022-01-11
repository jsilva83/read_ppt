# Import external modules.
from pptx import Presentation


class PPTXRead:

    def __init__(self, in_file_name):
        """"""
        self.presentation = Presentation(in_file_name)
        return

    def get_table_values(self, in_slide_nr) -> dict:
        """"""
        out_dict = {}
        # Get the pointer to slide 'in_slide_nr'.
        a_slide = self.presentation.slides[in_slide_nr]
        # Get a list of tables in slide 'in_slide_nr'.
        a_slide_shape_tables = []
        a_slide_shape_tables_titles = []
        for a_shape in a_slide.shapes:
            # if shape_type == 19, it means that it is a table.
            if a_shape.shape_type == 19:
                a_slide_shape_tables.append(a_shape.table)
            if a_shape.shape_type == 17:
                a_slide_shape_tables_titles.append(a_shape.text)
        # Create the output.
        # For each table create an array.
        # {'Inbound/Outbound': {'SD-03': '6.6%', 'FI-02': '84.4%', 'FI-02': '93.9'}, 'Outbound': {'EWM-02': '0.0%'}}
        a_n = 0
        for a_table in a_slide_shape_tables:
            a_n += 1
            a_title = a_slide_shape_tables_titles[a_n]
            out_dict[a_title] = {}
            # a_j is the row index.
            a_j = 0
            for a_row in a_table.rows:
                # If first row, skip.
                if a_j == 0:
                    a_j += 1
                    continue
                # a_k is the column index.
                a_k = 0
                for a_column in a_row.cells:
                    a_k += 1
                    # Get the object column and the percentage column.
                    if a_k == 1:
                        a_key = a_column.text_frame.text
                        a_key = a_key.replace(u'\xa0', u' ')
                    elif a_k == 4:
                        a_value = a_column.text_frame.text
                        a_value = float(a_value.replace(u'%', u''))
                    elif a_k > 4:
                        break
                out_dict[a_title][a_key] = a_value
        return out_dict


if __name__ == '__main__':
    # Test class.
    a_pptx = PPTXRead('10JAN2022 - HUF Mexico EDI Status Dashboard.pptx')
    slide_1_data_dict = a_pptx.get_table_values(in_slide_nr=1)
    slide_2_data_dict = a_pptx.get_table_values(in_slide_nr=2)
    pass
